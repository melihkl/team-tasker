from typing import Optional

from fastapi import FastAPI, Depends, HTTPException, status, Request, Form
from fastapi.responses import HTMLResponse, RedirectResponse
from fastapi.security import OAuth2PasswordBearer, OAuth2PasswordRequestForm
from fastapi.templating import Jinja2Templates
from fastapi.staticfiles import StaticFiles
from passlib.context import CryptContext
from jose import JWTError, jwt
from sqlalchemy import create_engine, Column, Integer, String, Date, ForeignKey, Enum
from sqlalchemy.ext.declarative import declarative_base
from sqlalchemy.orm import sessionmaker, relationship, Session
import enum
from datetime import datetime, timedelta
import os
from fastapi.responses import FileResponse
from pptx import Presentation
from pptx.util import Inches, Pt
from collections import defaultdict
import tempfile


# ----- Configuration -----
SECRET_KEY = "secret"
ALGORITHM = "HS256"
ACCESS_TOKEN_EXPIRE_MINUTES = 30

pwd_context = CryptContext(schemes=["bcrypt"], deprecated="auto")
oauth2_scheme = OAuth2PasswordBearer(tokenUrl="token")

app = FastAPI()
templates = Jinja2Templates(directory="templates")
app.mount("/static", StaticFiles(directory="static"), name="static")

if not os.path.exists("./templates"):
    os.makedirs("./templates")
if not os.path.exists("./static"):
    os.makedirs("./static")

DATABASE_URL = "sqlite:///./tasks.db"
engine = create_engine(DATABASE_URL, connect_args={"check_same_thread": False})
SessionLocal = sessionmaker(autocommit=False, autoflush=False, bind=engine)
Base = declarative_base()

# ----- Enums -----


class Role(str, enum.Enum):
    admin = "admin"
    user = "user"


class Status(str, enum.Enum):
    planned = "planned"
    ongoing = "ongoing"
    completed = "completed"
    canceled = "canceled"


# ----- Models -----
class User(Base):
    __tablename__ = "users"

    id = Column(Integer, primary_key=True, index=True)
    username = Column(String, unique=True, index=True)
    hashed_password = Column(String)
    role = Column(Enum(Role), default=Role.user)
    tasks = relationship("Task", back_populates="owner")


class Task(Base):
    __tablename__ = "tasks"

    id = Column(Integer, primary_key=True, index=True)
    project_name = Column(String)
    subject = Column(String)
    start_date = Column(Date)
    completion_date = Column(Date)
    status = Column(Enum(Status))
    description = Column(String)
    owner_id = Column(Integer, ForeignKey("users.id"))
    owner = relationship("User", back_populates="tasks")


class Risk(Base):
    __tablename__ = "risks"

    id = Column(Integer, primary_key=True, index=True)
    project_name = Column(String, nullable=False)
    description = Column(String, nullable=False)
    owner_id = Column(Integer, ForeignKey("users.id"))

    owner = relationship("User", back_populates="risks")


User.risks = relationship("Risk", back_populates="owner", cascade="all, delete")


Base.metadata.create_all(bind=engine)


# ----- Dependency -----
def get_db():
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()


# ----- Utility Functions -----
def get_password_hash(password):
    return pwd_context.hash(password)


def verify_password(plain_password, hashed_password):
    return pwd_context.verify(plain_password, hashed_password)


def create_access_token(data: dict, expires_delta: timedelta = None):
    to_encode = data.copy()
    if expires_delta:
        expire = datetime.utcnow() + expires_delta
    else:
        expire = datetime.utcnow() + timedelta(minutes=15)
    to_encode.update({"exp": expire})
    return jwt.encode(to_encode, SECRET_KEY, algorithm=ALGORITHM)


def get_user(db: Session, username: str):
    return db.query(User).filter(User.username == username).first()


def authenticate_user(db: Session, username: str, password: str):
    user = get_user(db, username)
    if not user:
        return False
    if not verify_password(password, user.hashed_password):
        return False
    return user


async def get_current_user(request: Request, db: Session = Depends(get_db)):
    token = request.cookies.get("Authorization")
    if token:
        credentials_exception = HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Could not validate credentials",
            headers={"WWW-Authenticate": "Bearer"},
        )
        try:
            payload = jwt.decode(token.split(" ")[1], SECRET_KEY, algorithms=[ALGORITHM])  # Bearer <token>
            username: str = payload.get("sub")
            if username is None:
                raise credentials_exception
        except JWTError:
            raise credentials_exception
        user = get_user(db, username)
        if user is None:
            raise credentials_exception
        return user


async def get_current_active_user(current_user: User = Depends(get_current_user)):
    return current_user


# ----- Routes -----
@app.post("/token")
async def login_for_access_token(form_data: OAuth2PasswordRequestForm = Depends(), db: Session = Depends(get_db)):
    user = authenticate_user(db, form_data.username, form_data.password)
    if not user:
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Incorrect username or password",
            headers={"WWW-Authenticate": "Bearer"},
        )
    access_token = create_access_token(data={"sub": user.username})
    return {"access_token": access_token, "token_type": "bearer"}


@app.get("/login", response_class=HTMLResponse)
async def login_form(request: Request):
    return templates.TemplateResponse("login.html", {"request": request, "error": None})


@app.post("/login", response_class=HTMLResponse)
async def login_submit(request: Request, username: str = Form(...), password: str = Form(...),
                       db: Session = Depends(get_db)):
    # Kullanıcı doğrulama işlemi
    user = authenticate_user(db, username, password)
    if not user:
        # Hatalı giriş durumu
        return templates.TemplateResponse("login.html", {"request": request, "error": "Invalid credentials"})

    # Token oluşturuluyor
    access_token = create_access_token(data={"sub": user.username})

    # Token'ı HTTPOnly cookie olarak set ediyoruz
    response = RedirectResponse(url="/", status_code=302)
    response.set_cookie(key="Authorization", value=f"Bearer {access_token}", httponly=True)

    return response


@app.get("/register", response_class=HTMLResponse)
async def register_form(request: Request):
    return templates.TemplateResponse("register.html", {"request": request})


@app.post("/register")
async def register(username: str = Form(...), password: str = Form(...),
                   role: Role = Form(Role.user), db: Session = Depends(get_db)):
    if get_user(db, username):
        raise HTTPException(status_code=400, detail="Username already registered")
    user = User(username=username, hashed_password=get_password_hash(password), role=role)
    db.add(user)
    db.commit()
    db.refresh(user)
    return RedirectResponse(url="/login", status_code=302)


@app.get("/change-password", response_class=HTMLResponse)
async def change_password_form(request: Request):
    return templates.TemplateResponse("change_password.html", {"request": request})


@app.post("/change-password", response_class=HTMLResponse)
async def change_password(
    request: Request,
    username: str = Form(...),
    old_password: str = Form(...),
    new_password: str = Form(...),
    confirm_password: str = Form(...),
    db: Session = Depends(get_db)
):
    user = get_user(db, username)

    if not user:
        return templates.TemplateResponse("change_password.html", {
            "request": request,
            "error": "User not found"
        })

    if not verify_password(old_password, user.hashed_password):
        return templates.TemplateResponse("change_password.html", {
            "request": request,
            "error": "Old password is incorrect"
        })

    if new_password != confirm_password:
        return templates.TemplateResponse("change_password.html", {
            "request": request,
            "error": "New passwords do not match"
        })

    user.hashed_password = get_password_hash(new_password)
    db.commit()

    return templates.TemplateResponse("change_password.html", {
        "request": request,
        "success": "Password changed successfully"
    })


@app.get("/", response_class=HTMLResponse)
def index(
    request: Request,
    username: str = None,
    project_name: str = None,
    start_date: str = None,
    end_date: str = None,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_active_user)
):
    query = db.query(Task)
    if username:
        query = query.join(User).filter(User.username == username)
    if project_name:
        query = query.filter(Task.project_name.ilike(f"%{project_name}%"))
    if start_date:
        query = query.filter(Task.start_date >= start_date)
    if end_date:
        query = query.filter(Task.completion_date <= end_date)



    if current_user.role == Role.admin:
        risks = db.query(Risk).all()
        tasks = query.all()
    else:
        risks = db.query(Risk).filter(Risk.owner_id == current_user.id).all()
        tasks = query.filter(Task.owner_id == current_user.id).all()

    return templates.TemplateResponse("index.html", {
        "request": request,
        "tasks": tasks,
        "risks": risks,
        "user": current_user,
        "query": {
            "username": username,
            "project_name": project_name,
            "start_date": start_date,
            "end_date": end_date
        }
    })


@app.get("/add", response_class=HTMLResponse)
async def add_form(request: Request, current_user: User = Depends(get_current_active_user)):
    if current_user.role == Role.user:
        return templates.TemplateResponse("add.html", {"request": request, "Status": Status})
    raise HTTPException(status_code=403, detail="Admin does not have permission to add tasks")


@app.post("/add")
async def add_task(request: Request,
                   project_name: str = Form(...),
                   subject: str = Form(...),
                   start_date: str = Form(...),
                   completion_date: str = Form(...),
                   status: Status = Form(...),
                   description: str = Form(...),
                   current_user: User = Depends(get_current_active_user),
                   db: Session = Depends(get_db)):
    if current_user.role == Role.user:
        task = Task(
            project_name=project_name,
            subject=subject,
            start_date=datetime.strptime(start_date, "%Y-%m-%d"),
            completion_date=datetime.strptime(completion_date, "%Y-%m-%d"),
            status=status,
            description=description,
            owner_id=current_user.id
        )
        db.add(task)
        db.commit()
        return RedirectResponse(url="/", status_code=302)
    raise HTTPException(status_code=403, detail="Admin does not have permission to add tasks")


@app.post("/delete/{task_id}")
async def delete_task(task_id: int, current_user: User = Depends(get_current_active_user), db: Session = Depends(get_db)):
    task = db.query(Task).filter(Task.id == task_id, Task.owner_id == current_user.id).first()
    if not task:
        raise HTTPException(status_code=404, detail="Task not found or not authorized")
    db.delete(task)
    db.commit()
    return RedirectResponse(url="/", status_code=302)


@app.get("/edit/{task_id}", response_class=HTMLResponse)
async def edit_form(task_id: int, request: Request, current_user: User = Depends(get_current_active_user), db: Session = Depends(get_db)):
    task = db.query(Task).filter(Task.id == task_id, Task.owner_id == current_user.id).first()
    if not task:
        raise HTTPException(status_code=404, detail="Task not found or not authorized")
    return templates.TemplateResponse("edit.html", {"request": request, "task": task, "Status": Status})


@app.post("/edit/{task_id}")
async def edit_task(task_id: int,
                    project_name: str = Form(...),
                    subject: str = Form(...),
                    start_date: str = Form(...),
                    completion_date: str = Form(...),
                    status: Status = Form(...),
                    description: str = Form(...),
                    current_user: User = Depends(get_current_active_user),
                    db: Session = Depends(get_db)):
    task = db.query(Task).filter(Task.id == task_id, Task.owner_id == current_user.id).first()
    if not task:
        raise HTTPException(status_code=404, detail="Task not found or not authorized")
    task.project_name = project_name
    task.subject = subject
    task.start_date = datetime.strptime(start_date, "%Y-%m-%d")
    task.completion_date = datetime.strptime(completion_date, "%Y-%m-%d")
    task.status = status
    task.description = description
    db.commit()
    return RedirectResponse(url="/", status_code=302)


@app.get('/logout', response_class=HTMLResponse)
def protected_route(request: Request, current_user: User = Depends(get_current_active_user)):
    resp = RedirectResponse(url="/login", status_code=status.HTTP_302_FOUND)
    return resp


@app.get("/export", response_class=FileResponse)
async def export_to_pptx(
    request: Request,
    username: Optional[str] = None,
    project_name: Optional[str] = None,
    start_date: Optional[str] = None,
    end_date: Optional[str] = None,
    status: Optional[Status] = None,
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    if current_user.role != Role.admin:
        raise HTTPException(status_code=403, detail="Only admins can export tasks.")

    # Görevleri çek
    query = db.query(Task)
    if username:
        query = query.join(User).filter(User.username == username)
    if project_name:
        query = query.filter(Task.project_name.ilike(f"%{project_name}%"))
    if start_date:
        query = query.filter(Task.start_date >= start_date)
    if end_date:
        query = query.filter(Task.completion_date <= end_date)
    if status:
        query = query.filter(Task.status == status)

    tasks = query.all()

    # Görevlerden elde edilen proje adlarını al
    filtered_project_names = {task.project_name for task in tasks}

    # Riskleri filtrele
    risk_query = db.query(Risk)
    if filtered_project_names:
        risk_query = risk_query.filter(Risk.project_name.in_(filtered_project_names))

    risks = risk_query.all()

    # Görevleri grupla
    grouped_tasks = defaultdict(lambda: defaultdict(list))
    for task in tasks:
        task_status = task.status.strip().lower()
        if task_status not in ["planned", "ongoing", "completed", "canceled"]:
            continue
        grouped_tasks[task.project_name][task_status].append(task)

    # Riskleri projelere göre grupla
    grouped_risks = defaultdict(list)
    for risk in risks:
        grouped_risks[risk.project_name].append(risk)

    # PPTX dosyası oluştur
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[5]

    for project_name in sorted(set(grouped_tasks.keys())):
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        title.text = f"{project_name}"

        left = Inches(0.5)
        top = Inches(1.0)
        width = Inches(9)
        height = Inches(5.5)

        content = ""

        if grouped_tasks[project_name].get("ongoing"):
            content += "🔸 Devam Eden Görevler:\n"
            for t in grouped_tasks[project_name]["ongoing"]:
                content += f"- {t.completion_date} -> {t.subject}\n"
            content += "\n"

        if grouped_tasks[project_name].get("completed"):
            content += "✅ Tamamlanan Görevler:\n"
            for t in grouped_tasks[project_name]["completed"]:
                content += f"- {t.completion_date} -> {t.subject}\n"
            content += "\n"

        if grouped_risks.get(project_name):
            content += "⚠️ Riskler / Sorunlar:\n"
            for r in grouped_risks[project_name]:
                content += f"- {r.description}\n"
            content += "\n"

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = content or "No tasks or risks available."
        p.font.size = Pt(14)

    # PPTX dosyasını geçici dosyaya kaydet
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
        pptx_path = tmp.name
        prs.save(pptx_path)

    return FileResponse(
        pptx_path,
        media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
        filename="tasks_export.pptx"
    )



@app.get("/add-risk", response_class=HTMLResponse)
async def add_risk_form(request: Request, current_user: User = Depends(get_current_active_user)):
    if current_user.role != Role.user:
        raise HTTPException(status_code=403, detail="Only users can add risks.")
    return templates.TemplateResponse("add_risk.html", {"request": request})


@app.post("/add-risk")
async def add_risk(
    request: Request,
    project_name: str = Form(...),
    description: str = Form(...),
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_active_user)
):
    if current_user.role != Role.user:
        raise HTTPException(status_code=403, detail="Only users can add risks.")

    new_risk = Risk(
        project_name=project_name,
        description=description,
        owner_id=current_user.id
    )
    db.add(new_risk)
    db.commit()
    return RedirectResponse(url="/", status_code=303)


@app.post("/deleteRisk/{risk_id}")
async def delete_risk(risk_id: int, current_user: User = Depends(get_current_active_user), db: Session = Depends(get_db)):
    risk = db.query(Risk).filter(Risk.id == risk_id, Risk.owner_id == current_user.id).first()
    if not risk:
        raise HTTPException(status_code=404, detail="Risk not found or not authorized")
    db.delete(risk)
    db.commit()
    return RedirectResponse(url="/", status_code=302)


@app.get("/risks", response_class=HTMLResponse)
async def view_risks(request: Request, db: Session = Depends(get_db), current_user: User = Depends(get_current_user)):
    if current_user.role == "admin":
        risks = db.query(Risk).all()
    else:
        risks = db.query(Risk).filter(Risk.owner_id == current_user.id).all()
    return templates.TemplateResponse("risks.html", {"request": request, "risks": risks, "user": current_user})

